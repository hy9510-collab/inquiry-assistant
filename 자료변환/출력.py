# 자료변환 — 출력 생성기
# 정리된 문서 구조(dict) → hwpx / xlsx / pptx / pdf / docx 파일.
# hwpx는 기존 hwpx_리포트생성.py의 OWPML 엔진을 그대로 재사용한다(원본 하나, 출력 여러 개).
import sys, re
from pathlib import Path

BASE = Path(__file__).resolve().parent.parent
if str(BASE) not in sys.path:
    sys.path.insert(0, str(BASE))
import importlib
_hg = importlib.import_module("hwpx_리포트생성")   # write_hwpx, HEADER_XML, SEC_PR, esc

MALGUN = Path("C:/Windows/Fonts/malgun.ttf")
MALGUN_BD = Path("C:/Windows/Fonts/malgunbd.ttf")


def _safe(name):
    name = re.sub(r'[\\/:*?"<>|]+', "_", str(name)).strip() or "정리결과"
    return name[:80]


# ── hwpx ──────────────────────────────────────────────────────────────
def _p(text, char=0, para=0, sec_pr=""):
    _p.n += 1
    return (f'<hp:p id="{_p.n}" paraPrIDRef="{para}" styleIDRef="0" pageBreak="0" columnBreak="0" merged="0">'
            f'<hp:run charPrIDRef="{char}">{sec_pr}<hp:t>{_hg.esc(text)}</hp:t></hp:run></hp:p>')
_p.n = 0


def _hwpx_section(doc):
    _p.n = 0
    body = [
        _p(doc["제목"], char=1, para=1, sec_pr=_hg.SEC_PR),
        _p(f"경기도의회 문화체육관광위원회 ｜ 정책지원관 정리"
           + (f" ｜ {doc['출처']}" if doc.get("출처") else ""), char=3, para=1),
        _p("", 0, 0),
    ]
    if doc.get("요약"):
        body.append(_p("요약", 2, 2))
        body += [_p("· " + str(x), 0, 0) for x in doc["요약"]]
    for s in doc.get("섹션", []):
        if s.get("소제목"):
            body.append(_p(s["소제목"], 2, 2))
        for para in s.get("문단", []):
            body.append(_p(str(para), 0, 0))
        t = s.get("표")
        if t:
            if t.get("제목"):
                body.append(_p(f"〔{t['제목']}〕", 4, 0))
            if t.get("헤더"):
                body.append(_p("  |  ".join(str(h) for h in t["헤더"]), 4, 0))
            for row in t.get("행", []):
                body.append(_p("  |  ".join("" if c is None else str(c) for c in row), 0, 0))
        body.append(_p("", 0, 0))
    if doc.get("비고"):
        body.append(_p("비고: " + str(doc["비고"]), 3, 0))
    return ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<hs:sec xmlns:hs="http://www.hancom.co.kr/hwpml2011/section"'
            ' xmlns:hp="http://www.hancom.co.kr/hwpml2011/paragraph"'
            ' xmlns:hc="http://www.hancom.co.kr/hwpml2011/core">' + "".join(body) + "</hs:sec>")


def gen_hwpx(doc, out_path):
    _hg.write_hwpx(Path(out_path), _hwpx_section(doc), title=doc["제목"],
                   preview=doc["제목"] + "\n" + "\n".join(doc.get("요약", [])))
    return out_path


# ── xlsx ──────────────────────────────────────────────────────────────
def gen_xlsx(doc, out_path):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    wb = Workbook()
    navy = Font(color="FFFFFF", bold=True)
    fill = PatternFill("solid", fgColor="1F3864")
    wrap = Alignment(wrap_text=True, vertical="top")

    ws = wb.active
    ws.title = "요약"
    ws["A1"] = doc["제목"]
    ws["A1"].font = Font(bold=True, size=14, color="1F3864")
    r = 3
    for x in doc.get("요약", []):
        ws.cell(r, 1, "· " + str(x)).alignment = wrap
        r += 1
    ws.column_dimensions["A"].width = 90

    ti = 0
    for s in doc.get("섹션", []):
        t = s.get("표")
        if not t:
            continue
        ti += 1
        title = _safe(t.get("제목") or s.get("소제목") or f"표{ti}")[:28]
        sh = wb.create_sheet(f"{ti}.{title}"[:31])
        cr = 1
        headers = t.get("헤더") or []
        if headers:
            for ci, h in enumerate(headers, 1):
                c = sh.cell(cr, ci, str(h)); c.font = navy; c.fill = fill
            cr += 1
        for row in t.get("행", []):
            for ci, v in enumerate(row, 1):
                sh.cell(cr, ci, "" if v is None else str(v)).alignment = wrap
            cr += 1
        for ci in range(1, (len(headers) or 1) + 1):
            sh.column_dimensions[sh.cell(1, ci).column_letter].width = 22

    # 표가 하나도 없으면 섹션 문단을 표로 정리
    if ti == 0 and doc.get("섹션"):
        sh = wb.create_sheet("본문")
        sh.cell(1, 1, "소제목").font = navy; sh.cell(1, 1).fill = fill
        sh.cell(1, 2, "내용").font = navy; sh.cell(1, 2).fill = fill
        cr = 2
        for s in doc["섹션"]:
            sh.cell(cr, 1, s.get("소제목", "")).alignment = wrap
            sh.cell(cr, 2, "\n".join(str(p) for p in s.get("문단", []))).alignment = wrap
            cr += 1
        sh.column_dimensions["A"].width = 24; sh.column_dimensions["B"].width = 80
    wb.save(out_path)
    return out_path


# ── pptx ──────────────────────────────────────────────────────────────
def gen_pptx(doc, out_path):
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    prs = Presentation()
    NAVY = RGBColor(0x1F, 0x38, 0x64)

    # 표지
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = doc["제목"]
    sub = doc.get("출처") or "경기도의회 문화체육관광위원회 ｜ 정책지원관"
    if s.placeholders and len(s.placeholders) > 1:
        s.placeholders[1].text = sub

    # 요약
    if doc.get("요약"):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = "요약"
        tf = s.placeholders[1].text_frame
        tf.clear()
        for i, x in enumerate(doc["요약"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(x); p.font.size = Pt(18)

    # 섹션별 슬라이드
    for sec in doc.get("섹션", []):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = sec.get("소제목") or "내용"
        tf = s.placeholders[1].text_frame
        tf.clear()
        paras = sec.get("문단", []) or [""]
        for i, para in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(para); p.font.size = Pt(16)
        t = sec.get("표")
        if t and t.get("행"):
            _pptx_table(s, prs, t)
    prs.save(out_path)
    return out_path


def _pptx_table(slide, prs, t):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    headers = t.get("헤더") or []
    rows = t.get("행", [])
    ncol = len(headers) or (len(rows[0]) if rows else 1)
    nrow = len(rows) + (1 if headers else 0)
    top = Inches(4.2)
    gtbl = slide.shapes.add_table(max(nrow, 1), max(ncol, 1),
                                  Inches(0.5), top, Inches(9), Inches(0.4 * max(nrow, 1))).table
    r0 = 0
    if headers:
        for c, h in enumerate(headers):
            cell = gtbl.cell(0, c); cell.text = str(h)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1F, 0x38, 0x64)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
        r0 = 1
    for ri, row in enumerate(rows):
        for c in range(ncol):
            v = row[c] if c < len(row) else ""
            cell = gtbl.cell(ri + r0, c); cell.text = "" if v is None else str(v)
            cell.text_frame.paragraphs[0].font.size = Pt(10)


# ── pdf ───────────────────────────────────────────────────────────────
def gen_pdf(doc, out_path):
    from fpdf import FPDF
    pdf = FPDF(format="A4")
    pdf.set_auto_page_break(True, margin=15)
    pdf.add_page()
    pdf.add_font("malgun", "", str(MALGUN))
    bold = "malgun"
    if MALGUN_BD.exists():
        pdf.add_font("malgun", "B", str(MALGUN_BD)); bold = "malgun"
    W = pdf.w - pdf.l_margin - pdf.r_margin

    pdf.set_font(bold, "B" if MALGUN_BD.exists() else "", 18)
    pdf.set_text_color(0x1F, 0x38, 0x64)
    pdf.multi_cell(W, 10, doc["제목"])
    pdf.set_font("malgun", "", 9); pdf.set_text_color(0x66, 0x66, 0x66)
    pdf.multi_cell(W, 6, "경기도의회 문화체육관광위원회 ｜ 정책지원관 정리"
                   + (f" ｜ {doc['출처']}" if doc.get("출처") else ""))
    pdf.ln(2)

    def heading(txt):
        pdf.ln(2); pdf.set_font(bold, "B" if MALGUN_BD.exists() else "", 13)
        pdf.set_text_color(0x1F, 0x38, 0x64); pdf.multi_cell(W, 8, txt); pdf.set_text_color(0, 0, 0)

    def body(txt):
        pdf.set_font("malgun", "", 11); pdf.multi_cell(W, 6.5, txt)

    if doc.get("요약"):
        heading("요약")
        for x in doc["요약"]:
            body("· " + str(x))
    for s in doc.get("섹션", []):
        if s.get("소제목"):
            heading(s["소제목"])
        for para in s.get("문단", []):
            body(str(para))
        t = s.get("표")
        if t and t.get("행"):
            _pdf_table(pdf, t, W, bold)
    if doc.get("비고"):
        pdf.ln(2); pdf.set_font("malgun", "", 9); pdf.set_text_color(0x66, 0x66, 0x66)
        pdf.multi_cell(W, 5.5, "비고: " + str(doc["비고"]))
    pdf.output(str(out_path))
    return out_path


def _pdf_table(pdf, t, W, bold):
    from fpdf import FPDF
    headers = [str(h) for h in (t.get("헤더") or [])]
    rows = t.get("행", [])
    ncol = len(headers) or (len(rows[0]) if rows else 1)
    cw = W / max(ncol, 1)
    lh = 6
    pdf.ln(1)
    if t.get("제목"):
        pdf.set_font(bold, "B" if MALGUN_BD.exists() else "", 10); pdf.multi_cell(W, 6, f"〔{t['제목']}〕")
    if headers:
        pdf.set_font(bold, "B" if MALGUN_BD.exists() else "", 9)
        pdf.set_fill_color(0xE8, 0xED, 0xF7)
        _pdf_row(pdf, headers, cw, lh, fill=True)
    pdf.set_font("malgun", "", 9)
    for row in rows:
        cells = ["" if c is None else str(c) for c in row]
        cells += [""] * (ncol - len(cells))
        _pdf_row(pdf, cells[:ncol], cw, lh, fill=False)
    pdf.ln(1)


def _pdf_row(pdf, cells, cw, lh, fill):
    # 각 셀 줄 수를 계산해 행 높이를 맞춘 뒤 그린다.
    heights = []
    for c in cells:
        lines = pdf.multi_cell(cw, lh, c, dry_run=True, output="LINES")
        heights.append(max(1, len(lines)))
    rh = max(heights) * lh
    x0, y0 = pdf.get_x(), pdf.get_y()
    if y0 + rh > pdf.h - pdf.b_margin:
        pdf.add_page(); x0, y0 = pdf.get_x(), pdf.get_y()
    for c in cells:
        x, y = pdf.get_x(), pdf.get_y()
        pdf.rect(x, y, cw, rh)
        if fill:
            pdf.rect(x, y, cw, rh, style="F"); pdf.set_xy(x, y)
        pdf.multi_cell(cw, lh, c, border=0, align="L")
        pdf.set_xy(x + cw, y)
    pdf.set_xy(x0, y0 + rh)


# ── docx ──────────────────────────────────────────────────────────────
def gen_docx(doc, out_path):
    import docx
    from docx.shared import Pt, RGBColor
    d = docx.Document()
    h = d.add_heading(doc["제목"], level=0)
    sub = d.add_paragraph("경기도의회 문화체육관광위원회 ｜ 정책지원관 정리"
                          + (f" ｜ {doc['출처']}" if doc.get("출처") else ""))
    sub.runs[0].font.size = Pt(9); sub.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    if doc.get("요약"):
        d.add_heading("요약", level=1)
        for x in doc["요약"]:
            d.add_paragraph(str(x), style="List Bullet")
    for s in doc.get("섹션", []):
        if s.get("소제목"):
            d.add_heading(s["소제목"], level=1)
        for para in s.get("문단", []):
            d.add_paragraph(str(para))
        t = s.get("표")
        if t and t.get("행"):
            headers = t.get("헤더") or []
            ncol = len(headers) or len(t["행"][0])
            table = d.add_table(rows=0, cols=ncol); table.style = "Light Grid Accent 1"
            if headers:
                cells = table.add_row().cells
                for i, hh in enumerate(headers):
                    cells[i].text = str(hh)
            for row in t["행"]:
                cells = table.add_row().cells
                for i in range(ncol):
                    cells[i].text = "" if i >= len(row) or row[i] is None else str(row[i])
    if doc.get("비고"):
        p = d.add_paragraph("비고: " + str(doc["비고"])); p.runs[0].font.size = Pt(9)
    d.save(out_path)
    return out_path


GENERATORS = {"hwpx": gen_hwpx, "xlsx": gen_xlsx, "pptx": gen_pptx, "pdf": gen_pdf, "docx": gen_docx}
EXT = {"hwpx": ".hwpx", "xlsx": ".xlsx", "pptx": ".pptx", "pdf": ".pdf", "docx": ".docx"}


def generate(doc, fmt, out_dir):
    fmt = (fmt or doc.get("추천형식") or "hwpx").lower()
    if fmt not in GENERATORS:
        fmt = "hwpx"
    out_dir = Path(out_dir); out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / (_safe(doc["제목"]) + EXT[fmt])
    n = 1
    while out.exists():
        out = out_dir / f"{_safe(doc['제목'])}_{n}{EXT[fmt]}"; n += 1
    GENERATORS[fmt](doc, out)
    return out
