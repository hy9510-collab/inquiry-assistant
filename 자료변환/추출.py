# 자료변환 — 입력 추출기
# 여러 형식(파일·링크)에서 텍스트를 뽑아 Gemini에 넘길 원문을 만든다.
# 동영상·음성·이미지는 텍스트로 뽑지 않고 파일 그대로 Gemini에 보낸다(정리.py에서 처리).
import re, html, zipfile
from pathlib import Path

# 미디어(원본을 그대로 Gemini에 업로드) 확장자 → MIME
MEDIA_MIME = {
    ".mp4": "video/mp4", ".mov": "video/quicktime", ".webm": "video/webm",
    ".mp3": "audio/mpeg", ".wav": "audio/wav", ".m4a": "audio/mp4", ".ogg": "audio/ogg",
    ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".webp": "image/webp",
}
# 텍스트로 추출하는 문서 확장자
DOC_EXTS = {".txt", ".md", ".csv", ".json", ".hwpx", ".xlsx", ".xlsm", ".pptx", ".docx", ".pdf"}


YOUTUBE_RE = re.compile(r"(?:youtube\.com|youtu\.be)", re.I)
YOUTUBE_ID_RE = re.compile(r"(?:v=|/live/|/shorts/|/embed/|youtu\.be/)([A-Za-z0-9_-]{11})")


def is_youtube(url):
    return bool(YOUTUBE_RE.search(url or ""))


def youtube_id(url):
    m = YOUTUBE_ID_RE.search(url or "")
    return m.group(1) if m else None


def from_youtube(url):
    """유튜브 영상의 자막(스크립트)을 가져온다. 한국어 우선, 없으면 아무 언어."""
    from youtube_transcript_api import YouTubeTranscriptApi
    vid = youtube_id(url)
    if not vid:
        raise ValueError("유튜브 영상 주소에서 영상 ID를 찾지 못했습니다.")
    api = YouTubeTranscriptApi()
    tr = None
    try:
        tr = api.fetch(vid, languages=["ko", "ko-KR", "en"])
    except Exception:
        try:
            tr = next(iter(api.list(vid))).fetch()   # 사용 가능한 첫 자막
        except Exception:
            raise ValueError("이 영상에는 가져올 수 있는 자막이 없습니다. "
                             "자막 없는 영상은 파일로 저장해 19MB 이하로 잘라서 올려주세요.")
    text = " ".join(s.text for s in tr).strip()
    if not text:
        raise ValueError("자막이 비어 있습니다.")
    return f"[유튜브 영상 자막]\n[출처] {url}\n\n{text}"


def is_media(name):
    return Path(name).suffix.lower() in MEDIA_MIME


def mime_of(name):
    return MEDIA_MIME.get(Path(name).suffix.lower(), "application/octet-stream")


def _unescape(s):
    return html.unescape(s)


def from_text(data: bytes):
    for enc in ("utf-8", "cp949", "euc-kr"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", "replace")


def from_hwpx(data: bytes):
    """hwpx(OWPML zip)의 본문 XML에서 문단 텍스트를 뽑는다."""
    out = []
    with zipfile.ZipFile(_bytesio(data)) as z:
        names = sorted(n for n in z.namelist() if re.match(r"Contents/section\d+\.xml$", n))
        for n in names:
            xml = z.read(n).decode("utf-8", "replace")
            xml = re.sub(r"</hp:p>", "\n", xml)          # 문단 경계 → 줄바꿈
            xml = re.sub(r"<hp:tab[^>]*/>", "\t", xml)
            text = re.sub(r"<[^>]+>", "", xml)            # 나머지 태그 제거
            out.append(_unescape(text))
    return "\n".join(ln for ln in "\n".join(out).splitlines() if ln.strip())


def from_xlsx(data: bytes):
    from openpyxl import load_workbook
    wb = load_workbook(_bytesio(data), data_only=True, read_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"### 시트: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                parts.append("\t".join(cells))
    wb.close()
    return "\n".join(parts)


def from_pptx(data: bytes):
    from pptx import Presentation
    prs = Presentation(_bytesio(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"### 슬라이드 {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        parts.append(t)
            if shape.has_table:
                for r in shape.table.rows:
                    parts.append("\t".join(c.text for c in r.cells))
    return "\n".join(parts)


def from_docx(data: bytes):
    import docx
    d = docx.Document(_bytesio(data))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for t in d.tables:
        for row in t.rows:
            parts.append("\t".join(c.text for c in row.cells))
    return "\n".join(parts)


def from_pdf(data: bytes):
    from pypdf import PdfReader
    reader = PdfReader(_bytesio(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages).strip()


def from_url(url: str):
    """웹페이지 본문 텍스트를 뽑는다."""
    import requests
    from bs4 import BeautifulSoup
    headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) UijeongAssistant/1.0"}
    r = requests.get(url, headers=headers, timeout=20)
    r.raise_for_status()
    soup = BeautifulSoup(r.content, "html.parser")
    for tag in soup(["script", "style", "noscript", "header", "footer", "nav", "aside"]):
        tag.decompose()
    title = (soup.title.string or "").strip() if soup.title else ""
    body = soup.find("article") or soup.find("main") or soup.body or soup
    lines = [ln.strip() for ln in body.get_text("\n").splitlines() if ln.strip()]
    text = "\n".join(lines)
    return (f"[제목] {title}\n[출처] {url}\n\n" + text) if title else text


_EXTRACTORS = {
    ".txt": from_text, ".md": from_text, ".csv": from_text, ".json": from_text,
    ".hwpx": from_hwpx, ".xlsx": from_xlsx, ".xlsm": from_xlsx,
    ".pptx": from_pptx, ".docx": from_docx, ".pdf": from_pdf,
}


def extract_bytes(name: str, data: bytes) -> str:
    """파일명 확장자에 맞는 추출기로 텍스트를 뽑는다."""
    ext = Path(name).suffix.lower()
    fn = _EXTRACTORS.get(ext)
    if fn is None:
        return from_text(data)      # 알 수 없는 형식은 텍스트로 시도
    if fn in (from_text,):
        return fn(data)
    return fn(data)


def extract_path(path) -> str:
    p = Path(path)
    return extract_bytes(p.name, p.read_bytes())


def _bytesio(data):
    import io
    return io.BytesIO(data)
